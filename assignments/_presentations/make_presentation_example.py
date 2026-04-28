import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE  = RGBColor(0x4E, 0x79, 0xA7)
GRAY  = RGBColor(0x99, 0x99, 0x99)
DARK  = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF0, 0xF0, 0xF0)
BORD  = RGBColor(0xCC, 0xCC, 0xCC)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── Helpers (identical to template) ───────────────────────────────────────────
def line(slide, left, top, width, height, color=BLUE):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def box(slide, left, top, width, height, fill=LGRAY, border=BORD):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(0.5)

def txt(slide, text, left, top, width, height,
        size=16, bold=False, italic=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def label(slide, text, left, top):
    txt(slide, text.upper(), left, top, Inches(6), Inches(0.35),
        size=15, bold=True, color=BLUE)

def content(slide, text, left, top, width, height, size=15):
    """Real content in a light box (non-italic, dark text)."""
    box(slide, left, top, width, height)
    txt(slide, text, left + Inches(0.15), top + Inches(0.12),
        width - Inches(0.3), height - Inches(0.2),
        size=size, color=DARK)

def slide_header(slide, title, slide_num):
    line(slide, 0, 0, W, Pt(6))
    txt(slide, title,
        Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
        size=24, bold=True, color=DARK)
    txt(slide, f"Slide {slide_num} of 3",
        Inches(11.3), Inches(0.28), Inches(1.7), Inches(0.45),
        size=12, color=GRAY, align=PP_ALIGN.RIGHT)
    line(slide, Inches(0.6), Inches(0.95), Inches(12.1), Pt(1),
         color=RGBColor(0xDD, 0xDD, 0xDD))


# ── SLIDE 1 · Title ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)

line(slide, 0, Inches(3.2), W, Pt(3))

txt(slide, "SOC 106 · Spring 2026",
    Inches(0.7), Inches(1.1), Inches(11.9), Inches(0.5),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)
txt(slide, "Does Parental Education Shape\nChildren's Educational Attainment?",
    Inches(0.7), Inches(1.75), Inches(11.9), Inches(1.5),
    size=34, bold=True, color=DARK, align=PP_ALIGN.CENTER)
txt(slide, "Alex Rivera  ·  May 5, 2026",
    Inches(0.7), Inches(3.45), Inches(11.9), Inches(0.5),
    size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 2 · Motivation & Research Question ───────────────────────────────────
slide = prs.slides.add_slide(blank)
slide_header(slide, "Motivation & Research Question", 1)

# Left column
label(slide, "Why it matters", Inches(0.6), Inches(1.1))
content(slide,
    "Education is the primary mechanism through which socioeconomic advantage is "
    "passed across generations. We know parental SES matters — but do mothers and "
    "fathers matter equally? And equally for sons and daughters?",
    Inches(0.6), Inches(1.5), Inches(5.9), Inches(1.6))

label(slide, "Research question", Inches(0.6), Inches(3.25))
content(slide,
    "Does parental education predict children's educational attainment, and does "
    "this relationship differ by the gender of the parent and the child?",
    Inches(0.6), Inches(3.65), Inches(5.9), Inches(0.85), size=16)

label(slide, "Theory / prior work", Inches(0.6), Inches(4.65))
content(slide,
    "Status attainment theory (Blau & Duncan 1967) vs. gender-socialization theory "
    "(Ridgeway 2011) — children may model same-gender parents more strongly.",
    Inches(0.6), Inches(5.05), Inches(5.9), Inches(1.0))

# Right column
label(slide, "Hypotheses", Inches(7.1), Inches(1.1))

txt(slide, "H₀ (Null hypothesis)",
    Inches(7.1), Inches(1.55), Inches(5.7), Inches(0.4),
    size=14, bold=True, color=DARK)
content(slide,
    "No differential effect of parent gender — parental education predicts equally "
    "regardless of which parent or child gender.",
    Inches(7.1), Inches(1.95), Inches(5.7), Inches(0.85))

txt(slide, "H₁ — Patriarchal Model",
    Inches(7.1), Inches(2.95), Inches(5.7), Inches(0.4),
    size=14, bold=True, color=DARK)
content(slide,
    "Father's education predicts children's attainment more strongly than mother's, "
    "regardless of child gender — reflecting traditional male-breadwinner stratification.",
    Inches(7.1), Inches(3.35), Inches(5.7), Inches(1.0))

txt(slide, "H₂ — Gender-Matching Model",
    Inches(7.1), Inches(4.5), Inches(5.7), Inches(0.4),
    size=14, bold=True, color=DARK)
content(slide,
    "Parental influence is gender-matched: father's education matters more for sons, "
    "mother's for daughters. Expected direction: positive interaction between parent "
    "and child gender.",
    Inches(7.1), Inches(4.9), Inches(5.7), Inches(1.1))

line(slide, Inches(6.75), Inches(1.05), Pt(1), Inches(5.9),
     color=RGBColor(0xDD, 0xDD, 0xDD))


# ── SLIDE 3 · Data, Variables & Method ────────────────────────────────────────
slide = prs.slides.add_slide(blank)
slide_header(slide, "Data, Variables & Method", 2)

# Left: Data
label(slide, "Data", Inches(0.6), Inches(1.1))
content(slide,
    "Source: General Social Survey (GSS), attain subset\n"
    "Sample: U.S. adults; n ≈ 1,517 with complete cases\n"
    "Unit of analysis: Individual respondent",
    Inches(0.6), Inches(1.5), Inches(3.6), Inches(1.5))

# Left: Method
label(slide, "Method", Inches(0.6), Inches(3.15))
content(slide,
    "Model: OLS regression with interaction term\n\n"
    "educ ~ paeduc + maeduc + sex_f\n"
    "  + paeduc×sex_f + maeduc×sex_f + age + race_f\n\n"
    "Why OLS? Outcome is continuous (0–20 yrs); interaction terms directly test H₂.\n\n"
    "Data prep: Listwise deletion for missing parental education. Race recoded to "
    "White / Non-white.",
    Inches(0.6), Inches(3.55), Inches(3.6), Inches(2.55))

# Right: Variables table
label(slide, "Key Variables", Inches(4.7), Inches(1.1))

line(slide, Inches(4.7), Inches(1.55), Inches(8.1), Pt(1.5))
for text, left, w in [
    ("Role",        Inches(4.7),  Inches(2.4)),
    ("Variable",    Inches(7.1),  Inches(2.0)),
    ("Measurement", Inches(9.1),  Inches(3.7)),
]:
    txt(slide, text, left, Inches(1.6), w, Inches(0.35),
        size=12, bold=True, color=BLUE)
line(slide, Inches(4.7), Inches(1.95), Inches(8.1), Pt(0.5),
     color=RGBColor(0xDD, 0xDD, 0xDD))

rows = [
    ("Outcome (DV)",       "educ",        "Years of education, 0–20"),
    ("Key predictor (IV)", "paeduc",       "Father's years of education, 0–20"),
    ("Key predictor (IV)", "maeduc",       "Mother's years of education, 0–20"),
    ("Moderator",          "sex_f",        "Male / Female (ref = Male)"),
    ("Controls",           "age, race_f",  "Age in years; White / Non-white"),
]
for i, (role, var, meas) in enumerate(rows):
    rt = Inches(2.05) + i * Inches(1.05)
    if i % 2 == 1:
        box(slide, Inches(4.7), rt, Inches(8.1), Inches(1.05),
            fill=LGRAY, border=RGBColor(0xEE, 0xEE, 0xEE))
    for text, left, w, ital, clr in [
        (role, Inches(4.8),  Inches(2.2), False, DARK),
        (var,  Inches(7.1),  Inches(1.9), False, DARK),
        (meas, Inches(9.1),  Inches(3.6), True,  GRAY),
    ]:
        txt(slide, text, left, rt + Pt(6), w, Inches(0.95),
            size=14, italic=ital, color=clr)
    line(slide, Inches(4.7), rt + Inches(1.05), Inches(8.1), Pt(0.5),
         color=RGBColor(0xEE, 0xEE, 0xEE))

line(slide, Inches(4.4), Inches(1.05), Pt(1), Inches(5.9),
     color=RGBColor(0xDD, 0xDD, 0xDD))


# ── Build figure ───────────────────────────────────────────────────────────────
labels_fig = ["Father's educ\n(sons)", "Father's educ\n(daughters)",
              "Mother's educ\n(sons)", "Mother's educ\n(daughters)"]
coefs      = [0.44, 0.29, 0.24, 0.41]
cis_lo     = [0.37, 0.22, 0.17, 0.33]
cis_hi     = [0.51, 0.36, 0.31, 0.49]
pt_colors  = ["#4E79A7", "#4E79A7", "#E15759", "#E15759"]

fig, ax = plt.subplots(figsize=(6.2, 3.4))
fig.patch.set_facecolor("#F0F0F0")
ax.set_facecolor("#F0F0F0")

for i, (c, lo, hi, col) in enumerate(zip(coefs, cis_lo, cis_hi, pt_colors)):
    ax.plot([lo, hi], [i, i], color=col, lw=2.5, zorder=2)
    ax.scatter(c, i, color=col, s=100, zorder=3)
    ax.text(c + 0.015, i + 0.15, "***", fontsize=10, color="#333333",
            fontfamily="Arial")

ax.axvline(0, color="#888888", lw=1, ls="--")
ax.set_yticks(range(len(labels_fig)))
ax.set_yticklabels(labels_fig, fontsize=11, fontfamily="Arial")
ax.set_xlabel("OLS Coefficient (years of education)", fontsize=11,
              fontfamily="Arial")
ax.set_title("Effect of Parental Education on Child's Education\nby Parent & Child Gender",
             fontsize=12, fontweight="bold", fontfamily="Arial")

ax.legend(handles=[
    mpatches.Patch(color="#4E79A7", label="Father's education"),
    mpatches.Patch(color="#E15759", label="Mother's education"),
], loc="lower right", fontsize=10, framealpha=0.9)

ax.set_xlim(-0.05, 0.65)
ax.spines["top"].set_visible(False)
ax.spines["right"].set_visible(False)
ax.grid(axis="x", alpha=0.3)
plt.tight_layout()

img_buf = io.BytesIO()
fig.savefig(img_buf, format="png", dpi=180, bbox_inches="tight")
img_buf.seek(0)
plt.close(fig)


# ── SLIDE 4 · Results & Takeaways ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
slide_header(slide, "Results & Takeaways", 3)

# Figure
label(slide, "Main figure", Inches(0.6), Inches(1.1))
slide.shapes.add_picture(img_buf, Inches(0.6), Inches(1.5),
                         Inches(7.3), Inches(3.9))

# Right: findings
label(slide, "Key findings", Inches(8.3), Inches(1.1))

txt(slide, "H₀ rejected",
    Inches(8.3), Inches(1.55), Inches(4.65), Inches(0.35),
    size=13, bold=True, color=DARK)
content(slide,
    "Parental gender does matter — effects are not symmetric across parents.",
    Inches(8.3), Inches(1.9), Inches(4.65), Inches(0.75), size=13)

txt(slide, "H₁ not supported",
    Inches(8.3), Inches(2.75), Inches(4.65), Inches(0.35),
    size=13, bold=True, color=DARK)
content(slide,
    "Father's education is not uniformly stronger — it predicts significantly "
    "more for sons (β = 0.44) than daughters (β = 0.29).",
    Inches(8.3), Inches(3.1), Inches(4.65), Inches(0.95), size=13)

txt(slide, "H₂ partially supported",
    Inches(8.3), Inches(4.15), Inches(4.65), Inches(0.35),
    size=13, bold=True, color=DARK)
content(slide,
    "Gender-matched pattern holds: mother's educ stronger for daughters "
    "(β = 0.41 vs. 0.24 for sons). All differences p < .001.",
    Inches(8.3), Inches(4.5), Inches(4.65), Inches(0.95), size=13)

line(slide, Inches(8.0), Inches(1.05), Pt(1), Inches(4.4),
     color=RGBColor(0xDD, 0xDD, 0xDD))

# Bottom strip
line(slide, Inches(0.6), Inches(5.55), Inches(12.1), Pt(1),
     color=RGBColor(0xDD, 0xDD, 0xDD))

for text, lbl, left, w in [
    ("Same-gender parents are stronger educational role models — consistent "
     "with gender-socialization theory.",
     "Main takeaway",      Inches(0.6),  Inches(3.85)),
    ("Cross-sectional data; no causal claim. Missing parental education may "
     "be non-random and likely underestimates effects.",
     "Limitations",        Inches(4.65), Inches(3.85)),
    ("Should I run separate models by gender rather than one interaction model?",
     "Question for class", Inches(8.7),  Inches(4.25)),
]:
    txt(slide, lbl.upper(), left, Inches(5.65), w, Inches(0.3),
        size=10, bold=True, color=BLUE)
    content(slide, text, left, Inches(6.0), w, Inches(1.25), size=13)


prs.save("final-presentation-example.pptx")
print("Saved: final-presentation-example.pptx")
