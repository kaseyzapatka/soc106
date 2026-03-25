from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE  = RGBColor(0x4E, 0x79, 0xA7)
GRAY  = RGBColor(0x99, 0x99, 0x99)
DARK  = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF0, 0xF0, 0xF0)
BORD  = RGBColor(0xCC, 0xCC, 0xCC)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def line(slide, left, top, width, height, color=BLUE):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def box(slide, left, top, width, height, fill=LGRAY, border=BORD):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(0.5)

def txt(slide, text, left, top, width, height,
        size=16, bold=False, italic=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def label(slide, text, left, top):
    """Small ALL-CAPS section label in blue."""
    txt(slide, text.upper(), left, top, Inches(6), Inches(0.35),
        size=15, bold=True, color=BLUE)

def placeholder(slide, text, left, top, width, height, size=15):
    """Italic gray hint text inside a light box."""
    box(slide, left, top, width, height)
    txt(slide, text, left + Inches(0.15), top + Inches(0.12),
        width - Inches(0.3), height - Inches(0.2),
        size=size, italic=True, color=GRAY)



# ── SLIDE 1 · Title ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)

line(slide, 0, Inches(3.2), W, Pt(3))                          # full-width rule

txt(slide, "SOC 106 · Spring 2026",
    Inches(0.7), Inches(1.1), Inches(11.9), Inches(0.5),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

txt(slide, "Your Research Paper Title",
    Inches(0.7), Inches(1.75), Inches(11.9), Inches(1.3),
    size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)

txt(slide, "Your Name  ·  May 5 or May 12, 2026",
    Inches(0.7), Inches(3.45), Inches(11.9), Inches(0.5),
    size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 2 · Motivation & Research Question ───────────────────────────────────
slide = prs.slides.add_slide(blank)

line(slide, 0, 0, W, Pt(6))                                     # thin top rule

txt(slide, "Motivation & Research Question",
    Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
    size=24, bold=True, color=DARK)
txt(slide, "Slide 1 of 3",
    Inches(11.3), Inches(0.28), Inches(1.7), Inches(0.45),
    size=12, color=GRAY, align=PP_ALIGN.RIGHT)

line(slide, Inches(0.6), Inches(0.95), Inches(12.1), Pt(1),
     color=RGBColor(0xDD, 0xDD, 0xDD))                          # thin divider

# Left column
label(slide, "Why it matters", Inches(0.6), Inches(1.1))
placeholder(slide,
    "What social problem or gap in knowledge motivated this study?\nWhy should the class care? Who is affected?",
    Inches(0.6), Inches(1.5), Inches(5.9), Inches(1.6))

label(slide, "Research question", Inches(0.6), Inches(3.25))
placeholder(slide,
    "State your research question in one clear sentence.",
    Inches(0.6), Inches(3.65), Inches(5.9), Inches(0.85),
    size=16)

label(slide, "Theory / prior work", Inches(0.6), Inches(4.65))
placeholder(slide,
    "Connect to one concept or prior finding that shapes your expectations.",
    Inches(0.6), Inches(5.05), Inches(5.9), Inches(1.0))

# Right column
label(slide, "Hypotheses", Inches(7.1), Inches(1.1))

txt(slide, "H₀ (Null hypothesis)",
    Inches(7.1), Inches(1.55), Inches(5.7), Inches(0.4),
    size=14, bold=True, color=DARK)
placeholder(slide,
    "No relationship / no difference between groups.",
    Inches(7.1), Inches(1.95), Inches(5.7), Inches(0.85))

txt(slide, "H₁ (Research hypothesis)",
    Inches(7.1), Inches(2.95), Inches(5.7), Inches(0.4),
    size=14, bold=True, color=DARK)
placeholder(slide,
    "What do you expect to find, and why?",
    Inches(7.1), Inches(3.35), Inches(5.7), Inches(1.0))

txt(slide, "Expected direction",
    Inches(7.1), Inches(4.5), Inches(5.7), Inches(0.4),
    size=14, bold=True, color=DARK)
placeholder(slide,
    "Positive / negative / no expected direction — and why.",
    Inches(7.1), Inches(4.9), Inches(5.7), Inches(0.75))

# vertical divider
line(slide, Inches(6.75), Inches(1.05), Pt(1), Inches(5.9),
     color=RGBColor(0xDD, 0xDD, 0xDD))


# ── SLIDE 3 · Data, Variables & Method ────────────────────────────────────────
slide = prs.slides.add_slide(blank)

line(slide, 0, 0, W, Pt(6))

txt(slide, "Data, Variables & Method",
    Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
    size=24, bold=True, color=DARK)
txt(slide, "Slide 2 of 3",
    Inches(11.3), Inches(0.28), Inches(1.7), Inches(0.45),
    size=12, color=GRAY, align=PP_ALIGN.RIGHT)

line(slide, Inches(0.6), Inches(0.95), Inches(12.1), Pt(1),
     color=RGBColor(0xDD, 0xDD, 0xDD))

# Left column: Data + Method
label(slide, "Data", Inches(0.6), Inches(1.1))
placeholder(slide,
    "Source: [Dataset name]\nSample: [Who is included? n = ?]\nUnit of analysis: [Individual / household / etc.]",
    Inches(0.6), Inches(1.5), Inches(3.6), Inches(1.5))

label(slide, "Method", Inches(0.6), Inches(3.15))
placeholder(slide,
    "Model: [OLS / Logistic / Chi-squared / t-test]\n\n"
    "Why this method? [One sentence.]\n\n"
    "Data decisions: [Recodes, dropped cases, restrictions.]",
    Inches(0.6), Inches(3.55), Inches(3.6), Inches(2.55))

# Right: variables table
label(slide, "Key Variables", Inches(4.7), Inches(1.1))

# Table header
line(slide, Inches(4.7), Inches(1.55), Inches(8.1), Pt(1.5))
for text, left, w in [
    ("Role",        Inches(4.7),  Inches(2.4)),
    ("Variable",    Inches(7.1),  Inches(2.0)),
    ("Measurement", Inches(9.1),  Inches(3.7)),
]:
    txt(slide, text, left, Inches(1.6), w, Inches(0.35),
        size=12, bold=True, color=BLUE)
line(slide, Inches(4.7), Inches(1.95), Inches(8.1), Pt(0.5),
     color=RGBColor(0xDD, 0xDD, 0xDD))

# Table rows
rows = [
    ("Outcome (DV)",       "[variable name]", "[How is it measured? e.g., income in $]"),
    ("Key predictor (IV)", "[variable name]", "[e.g., years of education, 0–20]"),
    ("Control 1",          "[variable name]", "[brief description]"),
    ("Control 2",          "[variable name]", "[brief description]"),
]
for i, (role, var, meas) in enumerate(rows):
    rt = Inches(2.05) + i * Inches(1.15)
    if i % 2 == 1:
        box(slide, Inches(4.7), rt, Inches(8.1), Inches(1.15),
            fill=LGRAY, border=RGBColor(0xEE, 0xEE, 0xEE))
    for text, left, w, clr, ital in [
        (role, Inches(4.8),  Inches(2.2), DARK, False),
        (var,  Inches(7.1),  Inches(1.9), DARK, False),
        (meas, Inches(9.1),  Inches(3.6), GRAY, True),
    ]:
        txt(slide, text, left, rt + Pt(6), w, Inches(1.0),
            size=14, italic=ital, color=clr)
    line(slide, Inches(4.7), rt + Inches(1.15), Inches(8.1), Pt(0.5),
         color=RGBColor(0xEE, 0xEE, 0xEE))

# vertical divider
line(slide, Inches(4.4), Inches(1.05), Pt(1), Inches(5.9),
     color=RGBColor(0xDD, 0xDD, 0xDD))


# ── SLIDE 4 · Results & Takeaways ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank)

line(slide, 0, 0, W, Pt(6))

txt(slide, "Results & Takeaways",
    Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
    size=24, bold=True, color=DARK)
txt(slide, "Slide 3 of 3",
    Inches(11.3), Inches(0.28), Inches(1.7), Inches(0.45),
    size=12, color=GRAY, align=PP_ALIGN.RIGHT)

line(slide, Inches(0.6), Inches(0.95), Inches(12.1), Pt(1),
     color=RGBColor(0xDD, 0xDD, 0xDD))

# Left: figure placeholder
label(slide, "Main figure or table", Inches(0.6), Inches(1.1))
placeholder(slide,
    "Paste your chart, regression table, or cross-tab here.\n\n"
    "(Delete this box and insert your image or table.)",
    Inches(0.6), Inches(1.5), Inches(7.3), Inches(3.8),
    size=14)

# Right: findings
label(slide, "Key findings", Inches(8.3), Inches(1.1))

txt(slide, "Finding 1",
    Inches(8.3), Inches(1.55), Inches(4.65), Inches(0.35),
    size=13, bold=True, color=DARK)
placeholder(slide,
    "Direction + magnitude + significance.\ne.g., Each additional year of education is associated with a 0.3 SD increase in income (p < .05).",
    Inches(8.3), Inches(1.9), Inches(4.65), Inches(1.1), size=13)

txt(slide, "Finding 2",
    Inches(8.3), Inches(3.1), Inches(4.65), Inches(0.35),
    size=13, bold=True, color=DARK)
placeholder(slide,
    "Secondary finding, subgroup pattern, or null result.",
    Inches(8.3), Inches(3.45), Inches(4.65), Inches(0.85), size=13)

# vertical divider
line(slide, Inches(8.0), Inches(1.05), Pt(1), Inches(4.2),
     color=RGBColor(0xDD, 0xDD, 0xDD))

# Bottom strip: takeaway / limitations / questions
line(slide, Inches(0.6), Inches(5.45), Inches(12.1), Pt(1),
     color=RGBColor(0xDD, 0xDD, 0xDD))

for text, label_text, left, w in [
    ("One sentence: what did you find and what does it mean sociologically?",
     "Main takeaway",     Inches(0.6),  Inches(3.85)),
    ("1–2 key limitations and how you'd address them with more time or data.",
     "Limitations",       Inches(4.65), Inches(3.85)),
    ("What specific feedback are you looking for from the class?",
     "Question for class", Inches(8.7),  Inches(4.25)),
]:
    txt(slide, label_text.upper(), left, Inches(5.55), w, Inches(0.3),
        size=10, bold=True, color=BLUE)
    placeholder(slide, text, left, Inches(5.9), w, Inches(1.35), size=13)


prs.save("final-presentation-template.pptx")
print("Saved: final-presentation-template.pptx")
